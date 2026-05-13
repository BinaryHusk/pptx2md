# Copyright 2024 Liu Siyao
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

import logging
import os
import re
import shutil
import tempfile
import uuid
import zipfile
import zlib
from pathlib import Path

from pptx import Presentation

logger = logging.getLogger(__name__)


def remove_unreadable_zip_members(file_path):
    file_path = os.fspath(file_path)
    repaired_path = os.path.join(tempfile.gettempdir(), f'{Path(file_path).stem}_repaired_{uuid.uuid4().hex}.pptx')
    skipped = []

    try:
        with zipfile.ZipFile(file_path) as zin, zipfile.ZipFile(repaired_path, 'w', zipfile.ZIP_DEFLATED) as zout:
            for info in zin.infolist():
                try:
                    data = zin.read(info.filename)
                except (zipfile.BadZipFile, zlib.error) as err:
                    skipped.append((info.filename, err))
                    continue
                zout.writestr(info, data)
    except (zipfile.BadZipFile, zlib.error):
        try:
            os.unlink(repaired_path)
        except OSError:
            pass
        raise

    if not skipped:
        try:
            os.unlink(repaired_path)
        except OSError:
            pass
        raise zipfile.BadZipFile('No unreadable zip members found to remove.')

    return repaired_path, skipped


def fix_null_rels(file_path):
    file_path = os.fspath(file_path)
    temp_dir_name = tempfile.mkdtemp()
    try:
        shutil.unpack_archive(file_path, temp_dir_name, 'zip')
        rels = [
            os.path.join(dp, f)
            for dp, dn, filenames in os.walk(temp_dir_name)
            for f in filenames
            if os.path.splitext(f)[1] == '.rels'
        ]
        pat = re.compile(r'<\S*Relationship[^>]+Target\S*=\S*"NULL"[^>]*/>', re.I)
        for fn in rels:
            with open(fn, 'r+') as f:
                content = f.read()
                res = pat.search(content)
                if res is not None:
                    content = pat.sub('', content)
                    f.seek(0)
                    f.truncate()
                    f.write(content)
        archive_base = os.path.join(tempfile.gettempdir(), f'{Path(file_path).stem}_purged_{uuid.uuid4().hex}')
        archive_path = shutil.make_archive(archive_base, 'zip', temp_dir_name)
        tgt = f'{archive_base}.pptx'
        shutil.move(archive_path, tgt)
        return tgt
    finally:
        shutil.rmtree(temp_dir_name, ignore_errors=True)


def load_pptx(file_path: str) -> Presentation:
    if not os.path.exists(file_path):
        logger.error(f'source file {file_path} not exist!')
        logger.error(f'absolute path: {os.path.abspath(file_path)}')
        raise FileNotFoundError(file_path)
    try:
        prs = Presentation(file_path)
    except (zipfile.BadZipFile, zlib.error) as err:
        repaired_path = None
        try:
            repaired_path, skipped = remove_unreadable_zip_members(file_path)
            skipped_names = ', '.join(name for name, _ in skipped)
            logger.warning(f'Removed unreadable zip member(s) from {file_path}: {skipped_names}')
            prs = Presentation(repaired_path)
        except (zipfile.BadZipFile, zlib.error, ValueError) as repair_err:
            raise ValueError(f'Cannot read PPTX file {file_path}: invalid or corrupted zip archive.') from repair_err
        finally:
            if repaired_path:
                try:
                    os.unlink(repaired_path)
                except OSError:
                    logger.warning(f'failed to remove temporary repaired file {repaired_path}.')
    except KeyError as err:
        if len(err.args) > 0 and re.match(r'There is no item named .*NULL.* in the archive', str(err.args[0])):
            logger.info('corrupted links found, trying to purge...')
            res_path = None
            try:
                res_path = fix_null_rels(file_path)
                logger.info(f'created temporary purged file {res_path}.')
                prs = Presentation(res_path)
            except:
                logger.error(
                    'failed to purge corrupted links, you can report this at https://github.com/ssine/pptx2md/issues')
                raise err
            finally:
                if res_path:
                    try:
                        os.unlink(res_path)
                    except OSError:
                        logger.warning(f'failed to remove temporary purged file {res_path}.')
        else:
            logger.error('unknown error, you can report this at https://github.com/ssine/pptx2md/issues')
            raise err
    return prs


def prepare_titles(title_path: Path) -> dict[str, int]:
    titles: dict[str, int] = {}
    with open(title_path, 'r', encoding='utf8') as f:
        indent = -1
        for line in f.readlines():
            cnt = 0
            while line[cnt] == ' ':
                cnt += 1
            if cnt == 0:
                titles[line.strip()] = 1
            else:
                if indent == -1:
                    indent = cnt
                    titles[line.strip()] = 2
                else:
                    titles[line.strip()] = cnt // indent + 1
    return titles


def rgb_to_hex(rgb):
    r, g, b = rgb
    return f'#{r:02x}{g:02x}{b:02x}'
