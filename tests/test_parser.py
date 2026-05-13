from pathlib import Path
from types import SimpleNamespace
from zipfile import ZipFile

from PIL import Image
from pptx import Presentation
from pptx2md.parser import parse
from pptx2md.types import ConversionConfig
from pptx2md.utils import load_pptx


def _corrupt_zip_member(path: Path, member_name: str) -> None:
    with ZipFile(path) as zf:
        info = zf.getinfo(member_name)
        data_offset = info.header_offset + 30 + len(info.filename.encode()) + len(info.extra)

    with open(path, "r+b") as f:
        f.seek(data_offset)
        byte = f.read(1)
        f.seek(data_offset)
        f.write(bytes([byte[0] ^ 0xFF]))


def test_parse_skips_notes_slide_without_text_frame():
    slide = SimpleNamespace(
        shapes=[],
        has_notes_slide=True,
        notes_slide=SimpleNamespace(notes_text_frame=None),
    )
    prs = SimpleNamespace(slides=[slide])
    config = ConversionConfig(
        pptx_path=Path("deck.pptx"),
        output_path=Path("out.md"),
        image_dir=Path("img"),
    )

    result = parse(config, prs)

    assert len(result.slides) == 1
    assert result.slides[0].notes == []


def test_parse_skips_picture_with_missing_relationship(tmp_path):
    image_path = tmp_path / "image.png"
    Image.new("RGB", (1, 1), "red").save(image_path)

    source = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), 0, 0)
    prs.save(source)

    _corrupt_zip_member(source, "ppt/media/image1.png")
    repaired_prs = load_pptx(source)
    config = ConversionConfig(
        pptx_path=source,
        output_path=tmp_path / "out.md",
        image_dir=tmp_path / "img",
    )

    result = parse(config, repaired_prs)

    assert len(result.slides) == 1
    assert result.slides[0].elements == []
