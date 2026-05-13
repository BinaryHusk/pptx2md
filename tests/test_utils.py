from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile

import pytest
from pptx import Presentation

from pptx2md.utils import fix_null_rels, load_pptx


def _write_pptx_with_null_relationship(path: Path) -> None:
    with ZipFile(path, "w", ZIP_DEFLATED) as zf:
        zf.writestr(
            "ppt/_rels/presentation.xml.rels",
            '<Relationships><Relationship Id="rId1" Target="NULL"/></Relationships>',
        )
        zf.writestr("[Content_Types].xml", "<Types/>")


def _corrupt_zip_member(path: Path, member_name: str) -> None:
    with ZipFile(path) as zf:
        info = zf.getinfo(member_name)
        data_offset = info.header_offset + 30 + len(info.filename.encode()) + len(info.extra)

    with open(path, "r+b") as f:
        f.seek(data_offset)
        byte = f.read(1)
        f.seek(data_offset)
        f.write(bytes([byte[0] ^ 0xFF]))


def test_fix_null_rels_uses_unique_temporary_outputs(tmp_path):
    source = tmp_path / "deck.pptx"
    _write_pptx_with_null_relationship(source)

    first = Path(fix_null_rels(source))
    second = Path(fix_null_rels(source))

    try:
        assert first != second
        assert first.exists()
        assert second.exists()
        assert not (tmp_path / "deck_purged.pptx").exists()
    finally:
        first.unlink(missing_ok=True)
        second.unlink(missing_ok=True)


def test_load_pptx_reports_unreadable_archive(tmp_path):
    source = tmp_path / "broken.pptx"
    source.write_bytes(b"not a valid pptx package")

    with pytest.raises(ValueError, match="Cannot read PPTX file"):
        load_pptx(source)


def test_load_pptx_recovers_when_extra_zip_member_is_corrupted(tmp_path):
    source = tmp_path / "deck.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(source)

    with ZipFile(source, "a", ZIP_DEFLATED) as zf:
        zf.writestr("ppt/media/broken.bin", b"this member will be corrupted")
    _corrupt_zip_member(source, "ppt/media/broken.bin")

    loaded = load_pptx(source)

    assert len(loaded.slides) == 1
